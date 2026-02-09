import google.generativeai as genai
import fitz  # PyMuPDF
from pptx import Presentation
import os
import json
from dotenv import load_dotenv

load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise ValueError("API Key Missing")

genai.configure(api_key=api_key)

def extract_text(uploaded_file):
    # ... (Kode sama persis seperti sebelumnya) ...
    text = ""
    try:
        if uploaded_file.type == "application/pdf":
            doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            for page in doc:
                text += page.get_text() + "\n"
        elif uploaded_file.name.endswith(".pptx"):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif uploaded_file.type == "text/plain" or uploaded_file.name.endswith(".md"):
            text = str(uploaded_file.read(), "utf-8")
    except Exception as e:
        return f"Error: {str(e)}"
    return text

def generate_quiz_modern(context_text, config):
    # Config default kalau user males
    model_id = config.get('model_id', 'gemini-1.5-flash')
    
    model = genai.GenerativeModel(
        model_name=model_id,
        generation_config={"response_mime_type": "application/json", "temperature": 0.2}
    )

    # PROMPT BARU: Struktur Objek di dalam Array Opsi
    prompt = f"""
    ROLE: Academic Examiner.
    TASK: Create {config['q_count']} Multiple Choice Questions.
    CONTEXT: {context_text[:500000]}
    
    CRITICAL INSTRUCTION:
    1. The FIRST option (Index 0) MUST be the Correct Answer.
    2. The other 3-4 options MUST be Distractors (Wrong).
    3. You MUST provide a 'rationale' for EACH option (Why it is correct or why it is wrong).
    
    OUTPUT FORMAT (JSON Array):
    [
      {{
        "question": "Question text...",
        "options": [
          {{ "text": "(Correct Answer)", "rationale": "This is correct because..." }},
          {{ "text": "(Wrong Answer 1)", "rationale": "Incorrect. This refers to..." }},
          {{ "text": "(Wrong Answer 2)", "rationale": "Incorrect. This is actually..." }}
        ]
      }}
    ]
    """

    try:
        response = model.generate_content(prompt)
        return json.loads(response.text)
    except Exception as e:
        return {"error": str(e)}